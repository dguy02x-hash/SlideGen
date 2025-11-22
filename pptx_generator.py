"""
PowerPoint Theme Generator for PresPilot

Features:
- AI-generated custom themes from user prompts
- Predefined professional themes
- Alternating image placeholder positions (left, right, top, bottom)
- AI-written speaker notes that expand on slide content
- Image placeholders on all content slides (except title and thank you)
- Customizable colors, fonts, and layout styles
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Import grammar checking function from server
try:
    from server import proofread_slide_text
except ImportError:
    # Fallback if import fails
    def proofread_slide_text(text, max_tokens=500):
        return text

class ThemeGenerator:
    """
    Generate themed PowerPoint presentations with image placeholders

    Supports both predefined themes and AI-generated custom themes.
    All content slides include image placeholders in alternating positions.
    """
    
    THEMES = {
        "Business Black and Yellow": {
            "primary_color": RGBColor(255, 200, 0),  # Golden yellow
            "secondary_color": RGBColor(0, 0, 0),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(255, 200, 0),  # Golden yellow
            "background": RGBColor(0, 0, 0),
            "font": "Roboto Black",
            "title_font": "Roboto Black",
            "layouts": ["right", "center"]  # Right = image on right, center = centered bullets only
        },
        "Autumn Brown and Orange": {
            "primary_color": RGBColor(210, 105, 30),
            "secondary_color": RGBColor(255, 140, 0),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(255, 140, 0),
            "background": RGBColor(139, 69, 19),
            "font": "Georgia",
            "title_font": "Georgia",
            "layouts": ["left", "right", "bottom", "top"]
        },
        "Simplistic Red and White": {
            "primary_color": RGBColor(220, 20, 60),
            "secondary_color": RGBColor(255, 255, 255),
            "text_color": RGBColor(50, 50, 50),
            "accent_color": RGBColor(220, 20, 60),
            "background": RGBColor(255, 255, 255),
            "font": "Impact",
            "title_font": "Impact",
            "layouts": ["left", "right", "top", "bottom"]
        },
        "Nature Green": {
            "primary_color": RGBColor(34, 139, 34),
            "secondary_color": RGBColor(144, 238, 144),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(50, 205, 50),
            "background": RGBColor(34, 139, 34),
            "font": "Verdana",
            "title_font": "Verdana",
            "layouts": ["right", "left", "top", "bottom"]
        },
        "Elegant Black and Gray": {
            "primary_color": RGBColor(64, 64, 64),
            "secondary_color": RGBColor(192, 192, 192),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(192, 192, 192),
            "background": RGBColor(45, 45, 45),
            "font": "Garamond",
            "title_font": "Garamond",
            "layouts": ["left", "right", "top", "bottom"]
        },
        "Ocean Blue": {
            "primary_color": RGBColor(0, 119, 182),
            "secondary_color": RGBColor(173, 216, 230),
            "text_color": RGBColor(255, 255, 255),
            "accent_color": RGBColor(0, 191, 255),
            "background": RGBColor(0, 91, 150),
            "font": "Calibri",
            "title_font": "Calibri",
            "layouts": ["right", "left", "top", "bottom"]
        },
        "Sunset Orange": {
            "primary_color": RGBColor(255, 140, 0),
            "secondary_color": RGBColor(255, 99, 71),
            "text_color": RGBColor(50, 50, 50),  # Dark text for light background
            "accent_color": RGBColor(255, 85, 0),
            "background": RGBColor(255, 85, 0),  # Deep orange for title/thank you
            "content_background": RGBColor(255, 211, 179),  # Light peach for content slides
            "title_text_color": RGBColor(0, 0, 0),  # Black text for title slides
            "font": "Lobster",
            "title_font": "Lobster",
            "layouts": ["right", "left", "top", "bottom"]
        },
        "Minimalist Gray": {
            "primary_color": RGBColor(80, 80, 80),
            "secondary_color": RGBColor(120, 120, 120),
            "text_color": RGBColor(80, 80, 80),
            "accent_color": RGBColor(100, 100, 100),
            "background": RGBColor(220, 220, 220),
            "font": "Times New Roman",
            "title_font": "Times New Roman",
            "layouts": ["right", "left", "bottom", "top"]
        }
    }
    
    def __init__(self, theme_name="Business Black and Yellow", custom_style=None):
        """
        Initialize with a theme (predefined or custom AI-generated)

        Args:
            theme_name: Name of predefined theme (used if custom_style is None)
            custom_style: Dict with AI-generated style config (overrides theme_name)
        """
        if custom_style:
            # Use AI-generated custom style
            self.theme_name = custom_style.get('theme_name', 'Custom Theme')
            self.theme = self._convert_custom_style(custom_style)
            self.is_custom = True
        else:
            # Use predefined theme
            self.theme_name = theme_name
            if theme_name not in self.THEMES:
                raise ValueError(f"Theme '{theme_name}' not found")
            self.theme = self.THEMES[theme_name]
            self.is_custom = False

        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        self.layout_index = 0  # Track current layout for alternating layouts

        # No template files - always use programmatic generation
        self.template_path = None
        self.template_prs = None

    def _get_template_path(self):
        """Check if a template PPTX file exists for this theme"""
        import os
        template_dir = "theme-templates"
        template_file = os.path.join(template_dir, f"{self.theme_name}.pptx")
        if os.path.exists(template_file):
            return template_file
        return None

    def _load_template(self):
        """Load template PPTX file"""
        if self.template_path:
            try:
                return Presentation(self.template_path)
            except Exception as e:
                print(f"Warning: Could not load template {self.template_path}: {e}")
                return None
        return None

    def _clone_slide(self, template_slide):
        """Clone a slide from the template into the current presentation - using safer XML copying"""
        from copy import deepcopy

        # Add blank slide
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Use XML-level copying for better fidelity
        try:
            # Get the slide's XML tree
            src_slide_element = template_slide._element
            dest_slide_element = slide._element

            # Copy background - safer approach
            try:
                if hasattr(src_slide_element, 'cSld') and hasattr(src_slide_element.cSld, 'bg'):
                    # Copy background element if it exists
                    if src_slide_element.cSld.bg is not None:
                        dest_slide_element.cSld.bg = deepcopy(src_slide_element.cSld.bg)
            except Exception as e:
                print(f"Note: Could not copy slide background: {e}")

            # Copy all shape elements
            if hasattr(src_slide_element.cSld, 'spTree'):
                # Clear existing shapes except the first (which is required)
                while len(dest_slide_element.cSld.spTree) > 1:
                    dest_slide_element.cSld.spTree.remove(dest_slide_element.cSld.spTree[-1])

                # Copy shapes from source (skip the first required element)
                # Skip images - only copy text boxes and shapes
                for shape_elem in src_slide_element.cSld.spTree[1:]:
                    try:
                        # Check if this is an image (pic element)
                        if shape_elem.tag.endswith('}pic'):
                            # Skip images - create blank placeholder instead
                            continue

                        new_shape_elem = deepcopy(shape_elem)
                        dest_slide_element.cSld.spTree.append(new_shape_elem)
                    except Exception as e:
                        print(f"Note: Could not copy shape: {e}")

        except Exception as e:
            print(f"Warning: XML-level copy failed, falling back to shape-by-shape copy: {e}")
            # Fallback to old method if XML copying fails
            for shape in template_slide.shapes:
                self._copy_shape(shape, slide)

        return slide

    def _copy_shape(self, src_shape, dest_slide):
        """Copy a shape from template to destination slide"""
        try:
            # Copy shape with position and size
            if src_shape.shape_type == 1:  # Auto shape
                new_shape = dest_slide.shapes.add_shape(
                    src_shape.auto_shape_type,
                    src_shape.left, src_shape.top,
                    src_shape.width, src_shape.height
                )

                # Copy fill
                if src_shape.fill.type == 1:  # Solid
                    new_shape.fill.solid()
                    new_shape.fill.fore_color.rgb = src_shape.fill.fore_color.rgb
                elif src_shape.fill.type == 0:  # No fill
                    new_shape.fill.background()

                # Copy line
                if hasattr(src_shape.line, 'color'):
                    new_shape.line.color.rgb = src_shape.line.color.rgb
                    new_shape.line.width = src_shape.line.width
                else:
                    new_shape.line.fill.background()

                # Copy text if present
                if src_shape.has_text_frame:
                    self._copy_text_frame(src_shape.text_frame, new_shape.text_frame)

            elif src_shape.shape_type == 17:  # Text box
                new_shape = dest_slide.shapes.add_textbox(
                    src_shape.left, src_shape.top,
                    src_shape.width, src_shape.height
                )
                self._copy_text_frame(src_shape.text_frame, new_shape.text_frame)

        except Exception as e:
            print(f"Warning: Could not copy shape: {e}")

    def _copy_text_frame(self, src_tf, dest_tf):
        """Copy text frame properties"""
        try:
            # Copy text frame properties
            dest_tf.word_wrap = src_tf.word_wrap
            if hasattr(src_tf, 'vertical_anchor'):
                dest_tf.vertical_anchor = src_tf.vertical_anchor

            # Clear existing paragraphs
            for _ in range(len(dest_tf.paragraphs) - 1):
                p = dest_tf.paragraphs[0]._element
                p.getparent().remove(p)

            # Copy paragraphs
            for i, src_p in enumerate(src_tf.paragraphs):
                if i == 0:
                    dest_p = dest_tf.paragraphs[0]
                else:
                    dest_p = dest_tf.add_paragraph()

                # Copy paragraph properties
                dest_p.text = src_p.text
                dest_p.level = src_p.level
                dest_p.alignment = src_p.alignment

                # Copy font properties
                if src_p.runs:
                    src_run = src_p.runs[0]
                    dest_run = dest_p.runs[0] if dest_p.runs else None
                    if dest_run:
                        dest_run.font.name = src_run.font.name
                        dest_run.font.size = src_run.font.size
                        dest_run.font.bold = src_run.font.bold
                        dest_run.font.italic = src_run.font.italic
                        if hasattr(src_run.font.color, 'rgb'):
                            dest_run.font.color.rgb = src_run.font.color.rgb
        except Exception as e:
            print(f"Warning: Could not copy text frame: {e}")

    def _replace_text_in_slide(self, slide, replacements):
        """Replace placeholder text in a slide - intelligently handles all filler text"""
        # Common filler text patterns to replace
        common_fillers = {
            'Your Title Here': 'title',
            'Title Goes Here': 'title',
            'Presentation Title': 'title',
            'Title Here': 'title',
            'Enter Title': 'title',
            'Add Title': 'title',
            'Main Title': 'title',
            'Slide Title': 'title',
            'Your Name Here': 'presenter',
            'Presenter Name': 'presenter',
            'Speaker Name': 'presenter',
            'Author Name': 'presenter',
            'Your Name': 'presenter',
            'Name': 'presenter',
            'Section Title': 'section',
            'Topic': 'section',
            'Subtitle': 'section',
            'Add Text': 'section',
            'Content Here': 'section',
        }

        # Create comprehensive replacement map
        full_replacements = replacements.copy()

        # Map common fillers to their replacement values
        for filler, type_key in common_fillers.items():
            if type_key == 'title' and '{{TITLE}}' in replacements:
                full_replacements[filler] = replacements['{{TITLE}}']
            elif type_key == 'presenter' and '{{PRESENTER}}' in replacements:
                full_replacements[filler] = replacements['{{PRESENTER}}']
            elif type_key == 'section' and '{{SECTION_TITLE}}' in replacements:
                full_replacements[filler] = replacements['{{SECTION_TITLE}}']

        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame

                # Enable text wrapping and auto-size
                tf.word_wrap = True

                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text

                        # Replace all placeholders
                        for placeholder, value in full_replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, value)

                        # Text was replaced - word wrap will handle overflow

    def _replace_bullets_in_slide(self, slide, bullets):
        """Replace bullet points in a slide - intelligently detects and replaces all bullet formats"""
        if not bullets:
            return

        # Common bullet/list text patterns
        filler_patterns = [
            'bullet point', 'key point', 'main point', 'item', 'example',
            'feature', 'benefit', 'detail', 'point', 'add text', 'text here',
            'your text', 'content', 'lorem ipsum'
        ]

        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame

                # Enable text wrapping (NO auto-sizing - causes overflow)
                tf.word_wrap = True

                # Count how many paragraphs look like bullets
                bullet_paragraphs = []
                for i, p in enumerate(tf.paragraphs):
                    text = p.text.strip().lower()

                    # Check if it's a bullet point
                    is_bullet = (
                        p.text.startswith('•') or
                        p.text.startswith('-') or
                        p.text.startswith('*') or
                        p.text.startswith('→') or
                        p.text.startswith('►') or
                        p.level > 0 or  # Indented = bullet
                        any(pattern in text for pattern in filler_patterns) or
                        (len(text) > 0 and len(text) < 100 and not text.endswith(':'))  # Short text likely a bullet
                    )

                    if is_bullet:
                        bullet_paragraphs.append(i)

                # If we found bullet paragraphs, replace them
                if bullet_paragraphs and len(bullet_paragraphs) <= len(bullets):
                    for i, para_idx in enumerate(bullet_paragraphs):
                        if i < len(bullets):
                            p = tf.paragraphs[para_idx]
                            bullet_text = bullets[i]

                            # Preserve bullet character if present
                            if p.text.strip().startswith('•'):
                                p.text = f"• {bullet_text}"
                            elif p.text.strip().startswith('-'):
                                p.text = f"- {bullet_text}"
                            elif p.text.strip().startswith('*'):
                                p.text = f"* {bullet_text}"
                            elif p.text.strip().startswith('→'):
                                p.text = f"→ {bullet_text}"
                            elif p.text.strip().startswith('►'):
                                p.text = f"► {bullet_text}"
                            else:
                                p.text = bullet_text

                # If no bullets found but text frame has multiple paragraphs, try replacing them
                elif len(tf.paragraphs) >= len(bullets) and len(bullets) > 1:
                    for i, bullet_text in enumerate(bullets):
                        if i < len(tf.paragraphs):
                            tf.paragraphs[i].text = f"• {bullet_text}"

    def _ensure_text_fits(self, slide):
        """Ensure all text in slide fits properly - aggressive fitting for presentation mode"""
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame

                # CRITICAL: Enable word wrap everywhere
                tf.word_wrap = True

                # Set vertical anchor to top to prevent centering issues
                try:
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                except:
                    pass

                # Count number of bullet points and total text
                total_text = ''.join([p.text for p in tf.paragraphs])
                num_bullets = len([p for p in tf.paragraphs if p.text.strip()])

                # Calculate appropriate font size based on content amount
                # More aggressive sizing for presentation mode visibility
                if num_bullets > 4 or len(total_text) > 400:  # Lots of content
                    target_size = 13
                elif num_bullets > 3 or len(total_text) > 300:  # Medium content
                    target_size = 14
                elif len(total_text) > 200:  # Some content
                    target_size = 16
                else:  # Normal content
                    target_size = 18

                # Apply font size to all text
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        try:
                            # Only reduce font size, never increase
                            if run.font.size:
                                if run.font.size.pt > target_size:
                                    run.font.size = Pt(target_size)
                            else:
                                # If no font size set, set to target
                                run.font.size = Pt(target_size)
                        except:
                            pass

    def _add_background_image(self, slide, image_filename):
        """Add a background image from theme-templates folder"""
        template_dir = "theme-templates"
        image_path = os.path.join(template_dir, image_filename)

        if os.path.exists(image_path):
            # Add image as background - fills entire slide
            slide.shapes.add_picture(
                image_path,
                0, 0,
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
        else:
            print(f"Warning: Background image not found: {image_path}")

    def _convert_custom_style(self, custom_style):
        """Convert AI-generated style config to internal theme format"""
        def hex_to_rgb(hex_color):
            """Convert hex color to RGBColor and return tuple (RGBColor, r, g, b)"""
            hex_color = hex_color.lstrip('#')
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b), r, g, b

        # Determine image placeholder color based on style
        bg_color_hex = custom_style.get('background_color', '#FFFFFF')
        bg_rgb, bg_r, bg_g, bg_b = hex_to_rgb(bg_color_hex)
        # Calculate brightness using standard formula
        brightness = (bg_r * 299 + bg_g * 587 + bg_b * 114) / 1000

        # Choose placeholder color based on background brightness and style preference
        placeholder_style = custom_style.get('image_placeholder_style', 'dark')
        if placeholder_style == 'light':
            placeholder_color = RGBColor(240, 240, 240)
        elif placeholder_style == 'dark':
            placeholder_color = RGBColor(40, 40, 40)
        else:  # 'themed' - use a lighter/darker version of background
            if brightness > 128:  # Light background
                placeholder_color = RGBColor(
                    max(0, bg_r - 30),
                    max(0, bg_g - 30),
                    max(0, bg_b - 30)
                )
            else:  # Dark background
                placeholder_color = RGBColor(
                    min(255, bg_r + 30),
                    min(255, bg_g + 30),
                    min(255, bg_b + 30)
                )

        return {
            "primary_color": hex_to_rgb(custom_style.get('primary_color', '#FFD700'))[0],
            "secondary_color": hex_to_rgb(custom_style.get('secondary_color', '#000000'))[0],
            "text_color": hex_to_rgb(custom_style.get('text_color', '#FFFFFF'))[0],
            "accent_color": hex_to_rgb(custom_style.get('accent_color', '#FFD700'))[0],
            "background": hex_to_rgb(custom_style.get('background_color', '#000000'))[0],
            "font": custom_style.get('body_font', 'Arial'),
            "title_font": custom_style.get('title_font', 'Arial'),
            "title_size": custom_style.get('title_size', 36),
            "body_size": custom_style.get('body_size', 18),
            "placeholder_color": placeholder_color,
            "layouts": ["right", "left", "top", "bottom"]  # Image positions cycle through these
        }
    
    def _get_layout_for_slide(self):
        """Get layout position for current slide (cycles through layouts)"""
        layouts = self.theme["layouts"]
        return layouts[self.slide_count % len(layouts)]
    
    def add_title_slide(self, title, presenter_name="Your Name"):
        """Add a title slide - programmatically generated based on theme"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Theme-specific title slides matching PNG previews
        if self.theme_name == "Sunset Orange":
            # Orange gradient background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.gradient()
            bg.fill.gradient_angle = 90.0  # Vertical gradient
            # Orange to lighter orange gradient
            bg.fill.gradient_stops[0].color.rgb = RGBColor(255, 140, 0)  # Dark orange
            bg.fill.gradient_stops[1].color.rgb = RGBColor(255, 200, 100)  # Light orange
            bg.line.fill.background()

            # Title - centered, middle of slide, black text, Lobster font
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            p.alignment = PP_ALIGN.CENTER

            # Subtitle - centered below title, gray/brown text, Lobster font
            by_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.5), Inches(8), Inches(0.8)
            )
            tf = by_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Presented By {presenter_name}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(28)
            p.font.color.rgb = RGBColor(120, 90, 60)  # Gray/brown
            p.alignment = PP_ALIGN.CENTER

        elif self.theme_name == "Business Black and Yellow":
            # Load title background image
            self._add_background_image(slide, "Business Black and Yellow Title Background.jpg")

            # Title - white, centered
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Subtitle - white, centered below title
            by_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.2), Inches(8), Inches(0.7)
            )
            tf = by_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Presented by {presenter_name}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

        elif self.theme_name == "Ocean Blue":
            # Blue background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0, 120, 215)  # Ocean blue
            bg.line.fill.background()

            # Title box - rounded rectangle, top-left
            title_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(0.5), Inches(6.5), Inches(4.5)
            )
            title_box.fill.solid()
            title_box.fill.fore_color.rgb = RGBColor(0, 80, 150)  # Darker blue
            title_box.line.fill.background()

            # Title text
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Inches(0.3)
            tf.margin_top = Inches(0.3)
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(60)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

            # Presenter box - bottom-left
            pres_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.5), Inches(5.5), Inches(4), Inches(1.8)
            )
            pres_box.fill.solid()
            pres_box.fill.fore_color.rgb = RGBColor(0, 80, 150)
            pres_box.line.fill.background()

            tf = pres_box.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = f"Presented by [{presenter_name}]"
            p.font.name = self.theme["font"]
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

            # Small decorative box
            deco_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(4.8), Inches(5.5), Inches(1.9), Inches(1.8)
            )
            deco_box.fill.solid()
            deco_box.fill.fore_color.rgb = RGBColor(135, 206, 235)  # Light blue
            deco_box.line.fill.background()

        elif self.theme_name == "Simplistic Red and White":
            # Beige/Light gray background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(240, 240, 235)  # Light beige
            bg.line.fill.background()

            # Large diagonal red triangle from bottom-right
            triangle = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(6.5), Inches(3),
                Inches(4), Inches(4.5)
            )
            triangle.fill.solid()
            triangle.fill.fore_color.rgb = RGBColor(180, 30, 50)  # Dark red
            triangle.line.fill.background()
            triangle.rotation = 135  # Rotate to create diagonal effect

            # Title - large, red, on LEFT side
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2), Inches(5.5), Inches(2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(72)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent_color"]
            p.alignment = PP_ALIGN.LEFT

            # Presenter subtitle - smaller, red, below title
            by_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4), Inches(5.5), Inches(0.6)
            )
            tf = by_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"Presented by {presenter_name}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent_color"]
            p.alignment = PP_ALIGN.LEFT

            # Image placeholder box on RIGHT side
            img_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(7), Inches(1), Inches(3), Inches(5.5)
            )
            img_box.fill.solid()
            img_box.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Light gray
            img_box.line.color.rgb = RGBColor(100, 100, 100)  # Gray border
            img_box.line.width = Pt(2)

        elif self.theme_name == "Minimalist Gray":
            # Light gray background
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.theme["background"]
            bg.line.fill.background()

            # Title - centered at top
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(1)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_color"]
            p.alignment = PP_ALIGN.CENTER

            # Horizontal line under title
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.4), Inches(1.6), Inches(9.2), Pt(2)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = self.theme["text_color"]
            line.line.fill.background()

            # Subtitle - centered below title
            by_box = slide.shapes.add_textbox(
                Inches(1), Inches(1.8), Inches(8), Inches(0.5)
            )
            tf = by_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"[{presenter_name}]"
            p.font.name = self.theme["font"]
            p.font.size = Pt(22)
            p.font.italic = True
            p.font.color.rgb = self.theme["secondary_color"]
            p.alignment = PP_ALIGN.CENTER

            # Large image placeholder in center
            img_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(2.5), Inches(2.8), Inches(5), Inches(3.5)
            )
            img_box.fill.solid()
            img_box.fill.fore_color.rgb = RGBColor(140, 140, 140)
            img_box.line.fill.background()

        else:
            # Default/generic title slide for other themes
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.theme["background"]
            bg.line.fill.background()

            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(54)
            p.font.bold = True
            p.font.color.rgb = self.theme.get("title_text_color", self.theme["text_color"])
            p.alignment = PP_ALIGN.CENTER

            by_box = slide.shapes.add_textbox(
                Inches(1), Inches(4.8), Inches(8), Inches(0.6)
            )
            tf = by_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"By {presenter_name}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(24)
            p.font.italic = True
            p.font.color.rgb = self.theme.get("title_text_color", self.theme["secondary_color"])
            p.alignment = PP_ALIGN.CENTER

        return slide
    
    def add_content_slide(self, title, bullets, notes=""):
        """Add content slide - programmatically generated based on theme"""
        self.slide_count += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

        # Use theme-specific layouts that match PNG previews
        if self.theme_name == "Sunset Orange":
            self._add_sunset_orange_content(slide, title, bullets)
        elif self.theme_name == "Minimalist Gray":
            self._add_minimalist_gray_content(slide, title, bullets)
        elif self.theme_name == "Ocean Blue":
            self._add_ocean_blue_content(slide, title, bullets)
        elif self.theme_name == "Simplistic Red and White":
            self._add_simplistic_red_content(slide, title, bullets)
        elif self.theme_name == "Business Black and Yellow":
            self._add_business_black_content(slide, title, bullets)
        else:
            # Fallback for any other themes
            self._add_default_content(slide, title, bullets)

        # CRITICAL: Ensure all text fits properly in presentation mode
        self._ensure_text_fits(slide)

        return slide

    def _add_sunset_orange_content(self, slide, title, bullets):
        """Sunset Orange theme: Alternates between Body 1 and Body 2 layouts"""
        # Orange gradient background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.gradient()
        bg.fill.gradient_angle = 90.0  # Vertical gradient
        # Orange to lighter orange gradient
        bg.fill.gradient_stops[0].color.rgb = RGBColor(255, 140, 0)  # Dark orange
        bg.fill.gradient_stops[1].color.rgb = RGBColor(255, 200, 100)  # Light orange
        bg.line.fill.background()

        # Alternate between Body 1 and Body 2 layouts
        if self.slide_count % 2 == 1:
            # Body 1: Title top left, 3 bullets left, rounded rectangle image right

            # Title - top left, black, Lobster font
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(5), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["font"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black
            p.alignment = PP_ALIGN.LEFT

            # Bullets with dash markers - left side, gray/brown color
            y_pos = 2.0
            for i, bullet in enumerate(bullets[:3]):  # Max 3 bullets
                # Dash marker
                dash_box = slide.shapes.add_textbox(
                    Inches(0.6), Inches(y_pos), Inches(0.3), Inches(0.5)
                )
                tf = dash_box.text_frame
                p = tf.paragraphs[0]
                p.text = "-"
                p.font.name = self.theme["font"]
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(120, 90, 60)  # Gray/brown
                p.alignment = PP_ALIGN.LEFT

                # Bullet text
                bullet_box = slide.shapes.add_textbox(
                    Inches(1.1), Inches(y_pos), Inches(4.5), Inches(1.0)
                )
                tf = bullet_box.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.TOP
                p = tf.paragraphs[0]
                p.text = bullet
                p.font.name = self.theme["font"]
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(120, 90, 60)  # Gray/brown
                p.alignment = PP_ALIGN.LEFT

                y_pos += 1.3

            # Rounded rectangle image placeholder - right side
            img_placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6.0), Inches(2.0), Inches(3.5), Inches(4.5)
            )
            img_placeholder.fill.solid()
            img_placeholder.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Light gray
            img_placeholder.line.fill.background()

        else:
            # Body 2: Title top center, 2 circular images left (overlapping), 4 bullets right

            # Title - top center, black, Lobster font
            title_box = slide.shapes.add_textbox(
                Inches(2), Inches(0.5), Inches(6), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["font"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black
            p.alignment = PP_ALIGN.CENTER

            # Two circular image placeholders - left side, overlapping
            # First circle (back)
            circle1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.8), Inches(2.5), Inches(3.0), Inches(3.0)
            )
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = RGBColor(220, 220, 220)  # Light gray
            circle1.line.fill.background()

            # Second circle (front, overlapping)
            circle2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(2.2), Inches(3.0), Inches(3.0), Inches(3.0)
            )
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = RGBColor(200, 200, 200)  # Slightly darker gray
            circle2.line.fill.background()

            # Bullets with dash markers - right side, gray/brown color
            y_pos = 2.2
            for i, bullet in enumerate(bullets[:4]):  # Max 4 bullets
                # Dash marker
                dash_box = slide.shapes.add_textbox(
                    Inches(5.5), Inches(y_pos), Inches(0.3), Inches(0.5)
                )
                tf = dash_box.text_frame
                p = tf.paragraphs[0]
                p.text = "-"
                p.font.name = self.theme["font"]
                p.font.size = Pt(32)
                p.font.color.rgb = RGBColor(120, 90, 60)  # Gray/brown
                p.alignment = PP_ALIGN.LEFT

                # Bullet text
                bullet_box = slide.shapes.add_textbox(
                    Inches(6.0), Inches(y_pos), Inches(3.5), Inches(0.9)
                )
                tf = bullet_box.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.TOP
                p = tf.paragraphs[0]
                p.text = bullet
                p.font.name = self.theme["font"]
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(120, 90, 60)  # Gray/brown
                p.alignment = PP_ALIGN.LEFT

                y_pos += 1.0

    def _add_minimalist_gray_content(self, slide, title, bullets):
        """Minimalist Gray theme: Split layout with text left, image right"""
        # Light gray background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme["background"]
        bg.line.fill.background()

        # Title - italic with underline
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.8), Inches(5), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(48)
        p.font.italic = True
        p.font.color.rgb = self.theme["secondary_color"]
        p.alignment = PP_ALIGN.LEFT

        # Underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.65), Inches(5.5), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["secondary_color"]
        line.line.fill.background()

        # Bullet points on left
        y_pos = 2.3
        for i, bullet in enumerate(bullets[:3]):
            bullet_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_pos), Inches(4.5), Inches(0.6)
            )
            tf = bullet_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(24)
            p.font.color.rgb = self.theme["secondary_color"]
            p.alignment = PP_ALIGN.LEFT
            y_pos += 1.1

        # Tall image placeholder on right with gradient effect
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.2), Inches(0.8), Inches(3.3), Inches(6.2)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(140, 140, 140)
        img_placeholder.line.fill.background()

    def _add_ocean_blue_content(self, slide, title, bullets):
        """Ocean Blue theme: Single rounded rectangle box with uppercase bullets"""
        # Blue background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0, 120, 215)  # Ocean blue
        bg.line.fill.background()

        # Main content box - dark blue rounded rectangle
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(0.8), Inches(8.4), Inches(6)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = RGBColor(0, 80, 150)  # Darker blue
        content_box.line.fill.background()

        # Title inside box
        title_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.3), Inches(7.6), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

        # Uppercase bullet points
        y_pos = 3
        for i, bullet in enumerate(bullets[:3]):
            bullet_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(y_pos), Inches(7), Inches(0.8)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"• {bullet.upper()}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT
            y_pos += 1.3

    def _add_simplistic_red_content(self, slide, title, bullets):
        """Simplistic Red and White theme: Image left, bullets right, diagonal triangle bottom"""
        # AI Grammar check for Simplistic Red and White theme
        title = proofread_slide_text(title)
        bullets = [proofread_slide_text(bullet) for bullet in bullets]

        # Light beige background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(240, 240, 235)  # Light beige
        bg.line.fill.background()

        # Large diagonal red triangle at bottom (left to right)
        triangle = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE,
            Inches(0), Inches(4.5),
            Inches(10), Inches(3)
        )
        triangle.fill.solid()
        triangle.fill.fore_color.rgb = RGBColor(180, 30, 50)  # Dark red
        triangle.line.fill.background()
        triangle.rotation = 0  # Horizontal orientation

        # Title - large, red, centered at TOP
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(0.3), Inches(8), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.theme["accent_color"]
        p.alignment = PP_ALIGN.CENTER

        # Image placeholder on LEFT side
        img_placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.8), Inches(4), Inches(3.5)
        )
        img_placeholder.fill.solid()
        img_placeholder.fill.fore_color.rgb = RGBColor(250, 250, 250)  # Light gray
        img_placeholder.line.color.rgb = RGBColor(100, 100, 100)  # Gray border
        img_placeholder.line.width = Pt(2)

        # Red bullet points on RIGHT side
        y_pos = 2
        for i, bullet in enumerate(bullets[:4]):
            # Red bullet marker (dash)
            marker_box = slide.shapes.add_textbox(
                Inches(5.2), Inches(y_pos), Inches(0.3), Inches(0.4)
            )
            tf = marker_box.text_frame
            p = tf.paragraphs[0]
            p.text = "-"
            p.font.name = self.theme["font"]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent_color"]
            p.alignment = PP_ALIGN.CENTER

            # Bullet text
            bullet_box = slide.shapes.add_textbox(
                Inches(5.8), Inches(y_pos), Inches(4), Inches(0.8)
            )
            tf = bullet_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = bullet
            p.font.name = self.theme["font"]
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent_color"]
            p.alignment = PP_ALIGN.LEFT
            y_pos += 0.9

    def _add_business_black_content(self, slide, title, bullets):
        """Business Black and Yellow theme: Two layouts - right (with image) and center (no image)"""
        # Get current layout (cycles between 'right' and 'center')
        layout_type = self.theme["layouts"][self.layout_index % len(self.theme["layouts"])]

        # Load appropriate background image based on layout type
        if layout_type == "right":
            self._add_background_image(slide, "Business Black and Yellow Body 1 Background.jpg")
        else:  # center
            self._add_background_image(slide, "Business Black and Yellow Body 2 Background.jpg")

        if layout_type == "right":
            # Body 1 layout: Title top-left, bullets on left, image on right

            # Title - white, top-left
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(5), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

            # White bullet points on left
            y_pos = 2
            for i, bullet in enumerate(bullets[:3]):  # Only 3 bullets for this layout
                bullet_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(y_pos), Inches(5), Inches(0.8)
                )
                tf = bullet_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"- {bullet}"
                p.font.name = self.theme["font"]
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.LEFT
                y_pos += 1.2

            # Image placeholder on right
            img_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(6.5), Inches(1.8), Inches(3), Inches(3)
            )
            img_box.fill.solid()
            img_box.fill.fore_color.rgb = RGBColor(220, 220, 220)
            img_box.line.fill.background()

            # "INPUT IMAGE" text on placeholder
            img_text = slide.shapes.add_textbox(
                Inches(6.5), Inches(2.8), Inches(3), Inches(1)
            )
            tf = img_text.text_frame
            p = tf.paragraphs[0]
            p.text = "INPUT\nIMAGE"
            p.font.name = "Arial"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(100, 100, 100)
            p.alignment = PP_ALIGN.CENTER

        else:  # center layout
            # Body 2 layout: Title centered at top, bullets centered

            # Title - white, centered at top
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.8), Inches(8), Inches(0.8)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.theme["title_font"]
            p.font.size = Pt(52)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # White bullet points centered
            y_pos = 2.2
            for i, bullet in enumerate(bullets[:4]):  # Up to 4 bullets for this layout
                bullet_box = slide.shapes.add_textbox(
                    Inches(2.5), Inches(y_pos), Inches(5), Inches(0.8)
                )
                tf = bullet_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"- {bullet}"
                p.font.name = self.theme["font"]
                p.font.size = Pt(26)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
                y_pos += 1.1

        # Increment layout index for next slide
        self.layout_index += 1

    def _add_default_content(self, slide, title, bullets):
        """Fallback default layout for themes without specific layouts"""
        # Background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.theme.get("content_background", self.theme["background"])
        bg.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.theme.get("content_text_color", self.theme["text_color"])
        p.alignment = PP_ALIGN.LEFT

        # Bullets
        bullets_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(5.5), Inches(5)
        )
        tf = bullets_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.text = f"• {bullet}"
            p.font.name = self.theme["font"]
            p.font.size = Pt(18)
            p.font.color.rgb = self.theme.get("content_text_color", self.theme["text_color"])
            p.space_before = Pt(10)

        # Image placeholder
        img_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.5), Inches(1.8), Inches(3), Inches(5)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = RGBColor(100, 100, 100)
        img_box.line.color.rgb = self.theme.get("accent_color", self.theme["primary_color"])
        img_box.line.width = Pt(2)

    def add_thank_you_slide(self):
        """Add thank you slide - programmatically generated based on theme"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background - use background images for specific themes
        if self.theme_name == "Sunset Orange":
            # Gradient for Sunset Orange
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.gradient()
            bg.fill.gradient_angle = 90.0  # Vertical gradient
            # Orange to grey gradient
            bg.fill.gradient_stops[0].color.rgb = RGBColor(255, 140, 0)  # Orange
            bg.fill.gradient_stops[1].color.rgb = RGBColor(160, 160, 160)  # Grey
            bg.line.fill.background()
        elif self.theme_name == "Business Black and Yellow":
            # Load title background image (same as title slide)
            self._add_background_image(slide, "Business Black and Yellow Title Background.jpg")
        else:
            # Solid background for other themes
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0,
                self.prs.slide_width, self.prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = self.theme["background"]
            bg.line.fill.background()

        # "Thank You" text - centered
        thank_you_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        tf = thank_you_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(72)
        p.font.bold = True

        # Thank you text color - use title_text_color if available
        if self.is_custom:
            p.font.color.rgb = self.theme["primary_color"]
        elif "title_text_color" in self.theme:
            p.font.color.rgb = self.theme["title_text_color"]
        elif self.theme_name == "Simplistic Red and White":
            p.font.color.rgb = self.theme["accent_color"]
        else:
            p.font.color.rgb = self.theme["text_color"]

        p.alignment = PP_ALIGN.CENTER
        return slide

    def _add_slide_content(self, slide, title, bullets, layout):
        """Add title, content, and image placeholder based on layout"""

        # Use custom title size if available
        title_size = self.theme.get("title_size", 36)

        # Title bar for specific predefined themes only
        if not self.is_custom:
            if self.theme_name in ["Business Black and Yellow", "Autumn Brown and Orange"]:
                title_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    0, 0, self.prs.slide_width, Inches(1)
                )
                title_bar.fill.solid()
                title_bar.fill.fore_color.rgb = self.theme["primary_color"]
                title_bar.line.fill.background()

                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.15), Inches(9), Inches(0.7)
                )
                title_top = Inches(0.15)
            elif self.theme_name == "Simplistic Red and White":
                # Red accent bar on right
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(9.7), 0, Inches(0.3), self.prs.slide_height
                )
                accent.fill.solid()
                accent.fill.fore_color.rgb = self.theme["accent_color"]
                accent.line.fill.background()

                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(8.5), Inches(0.8)
                )
                title_top = Inches(0.5)
            else:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                )
                title_top = Inches(0.5)
        else:
            # Custom theme - simple title placement
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
            )
            title_top = Inches(0.5)

        # Title with overflow protection
        tf = title_box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.theme["title_font"]
        p.font.size = Pt(title_size)
        p.font.bold = True

        # Title color - use content_text_color if available (for themes with different content backgrounds)
        if self.is_custom:
            p.font.color.rgb = self.theme["primary_color"]
        elif "content_text_color" in self.theme:
            p.font.color.rgb = self.theme["content_text_color"]
        elif self.theme_name == "Business Black and Yellow":
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black on yellow
        elif self.theme_name == "Simplistic Red and White":
            p.font.color.rgb = self.theme["accent_color"]
        else:
            p.font.color.rgb = self.theme["text_color"]

        p.alignment = PP_ALIGN.LEFT
        
        # Add content and image based on layout
        if layout == "right":
            # Image on right, text on left
            self._add_image_placeholder(slide, Inches(6.5), Inches(1.8), Inches(3), Inches(5))
            self._add_bullet_text(slide, bullets, Inches(0.5), Inches(1.8), Inches(5.5), Inches(5))
        
        elif layout == "left":
            # Image on left, text on right
            self._add_image_placeholder(slide, Inches(0.5), Inches(1.8), Inches(3.5), Inches(5))
            self._add_bullet_text(slide, bullets, Inches(4.5), Inches(1.8), Inches(5), Inches(5))
        
        elif layout == "top":
            # Image on top, text below
            self._add_image_placeholder(slide, Inches(2), Inches(1.5), Inches(6), Inches(2.5))
            self._add_bullet_text(slide, bullets, Inches(0.5), Inches(4.3), Inches(9), Inches(2.8))
        
        elif layout == "bottom":
            # Text on top, image below
            self._add_bullet_text(slide, bullets, Inches(0.5), Inches(1.5), Inches(9), Inches(2.8))
            self._add_image_placeholder(slide, Inches(2), Inches(4.6), Inches(6), Inches(2.5))
    
    def _add_image_placeholder(self, slide, left, top, width, height):
        """Add a blank image placeholder box with border"""
        img_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        img_box.fill.solid()

        # Use custom placeholder color if available, otherwise use theme-specific defaults
        if self.is_custom and "placeholder_color" in self.theme:
            img_box.fill.fore_color.rgb = self.theme["placeholder_color"]
        elif self.theme_name == "Simplistic Red and White":
            img_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        elif self.theme_name == "Business Black and Yellow":
            img_box.fill.fore_color.rgb = RGBColor(40, 40, 40)
        else:
            img_box.fill.fore_color.rgb = RGBColor(60, 60, 60)

        # Border with accent color
        img_box.line.color.rgb = self.theme["accent_color"]
        img_box.line.width = Pt(2)
    
    def _add_bullet_text(self, slide, bullets, left, top, width, height):
        """Add bullet point text with overflow protection"""
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True

        # Enable auto-fit to shrink text if needed
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Use custom body size if available
        body_size = self.theme.get("body_size", 18)

        for i, bullet in enumerate(bullets):
            if i > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]

            p.text = bullet
            p.font.name = self.theme["font"]
            p.font.size = Pt(body_size)
            # Use content_text_color if available (for themes with different content backgrounds)
            p.font.color.rgb = self.theme.get("content_text_color", self.theme["text_color"])
            p.level = 0
            p.space_before = Pt(10)
    
    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)


def generate_presentation(title, topic, sections, theme_name="Business Black and Yellow",
                         notes_style="Detailed", slide_format="Detailed", custom_style=None, filename=None):
    """
    Generate a complete presentation with AI-written speaker notes

    Args:
        title: Presentation title
        topic: Main topic
        sections: List of dicts with 'title' and 'facts' keys
        theme_name: Name of predefined theme (used if custom_style is None)
        notes_style: Style of speaker notes (Concise, Detailed, Full Explanation)
        slide_format: Format of slide bullets (Concise = max 5 words, Detailed = full sentences)
        custom_style: Dict with AI-generated custom style (overrides theme_name)
        filename: Output filename

    Note: All content slides include image placeholders (except title and thank you slides)
    """
    if filename is None:
        filename = f"{title.replace(' ', '_')}.pptx"

    # Create generator with custom style or predefined theme
    gen = ThemeGenerator(theme_name=theme_name, custom_style=custom_style)

    # Add title slide
    gen.add_title_slide(title, "[Your Name]")

    # Add content slides
    for i, section in enumerate(sections):
        # Use pre-generated AI speaker notes if available, otherwise generate locally
        if 'speaker_notes' in section and section['speaker_notes']:
            notes = section['speaker_notes']
        else:
            # Fallback to local generation
            notes = generate_human_speaker_notes(
                section.get('title', f'Section {i+1}'),
                section.get('facts', []),
                section.get('notes_context', ''),
                notes_style,
                i + 1,
                section.get('custom_notes')
            )

        # Get facts - conversion should already be done in server.py before calling this
        facts = section.get('facts', ['Content for this section'])

        # Add content slide
        gen.add_content_slide(
            section.get('title', f'Section {i+1}'),
            facts,
            notes=notes
        )

    # Add thank you slide
    gen.add_thank_you_slide()

    gen.save(filename)
    return filename


def generate_human_speaker_notes(title, facts, context, style, slide_num, custom_notes=None):
    """Generate natural, human-sounding speaker notes that read like actual presentation speech"""

    facts_list = facts[:5] if len(facts) >= 5 else facts

    if style == "Concise":
        # Brief, natural explanation (4-6 sentences)
        notes = f"{title} represents an important aspect of this topic. "

        if facts_list:
            for i, fact in enumerate(facts_list):
                if i == 0:
                    notes += f"{fact}. "
                elif i == len(facts_list) - 1:
                    notes += f"Additionally, {fact}. "
                else:
                    notes += f"{fact}. "

        if context:
            notes += f"{context} "

        notes += "These elements work together to create a comprehensive understanding of the situation."

        return notes

    elif style == "Full Explanation":
        # In-depth, conversational explanation (10-14 sentences)
        notes = f"{title} encompasses several interconnected elements.\n\n"

        if context:
            notes += f"{context}\n\n"

        if facts_list:
            notes += f"{facts_list[0]}. This forms the foundation of our discussion today. "

            if len(facts_list) > 1:
                notes += f"Building on that, {facts_list[1]}. These two elements are closely related and support each other. "

            if len(facts_list) > 2:
                notes += f"Another significant factor involves {facts_list[2]}. This particular aspect has far-reaching implications. "

            if len(facts_list) > 3:
                notes += f"Furthermore, {facts_list[3]}. This adds considerable depth to our understanding. "

            if len(facts_list) > 4:
                notes += f"Finally, {facts_list[4]}. This completes the picture we've been building. "

        notes += f"\n\nThe significance of {title} becomes clear when we examine how these components interact. "
        notes += "Real-world applications demonstrate the practical value of understanding these relationships. "
        notes += "Organizations that successfully implement these principles often see measurable improvements in their outcomes."

        return notes

    else:  # Detailed (default)
        # Use custom notes if provided, otherwise format the facts
        if custom_notes:
            return custom_notes

        # Fallback: Format facts as a paragraph
        notes = f"{title}\n\n"

        if context:
            notes += f"{context}\n\n"

        if facts_list:
            # Join facts into a natural paragraph
            notes += ' '.join(facts_list)

        return notes


# Demo usage
if __name__ == "__main__":
    test_sections = [
        {
            "title": "Introduction to AI",
            "facts": [
                "Artificial intelligence transforms how we process information",
                "Machine learning enables systems to improve from experience",
                "Neural networks mimic human brain structure"
            ]
        },
        {
            "title": "Key Applications",
            "facts": [
                "Healthcare diagnostics with 95% accuracy",
                "Autonomous vehicles navigate complex environments",
                "Natural language processing powers chatbots"
            ]
        },
        {
            "title": "Future Implications",
            "facts": [
                "Ethical considerations guide AI development",
                "Job markets evolve with automation",
                "Human-AI collaboration enhances productivity"
            ]
        }
    ]
    
    # Test with different note styles
    for style in ["Concise", "Detailed", "Full Explanation"]:
        print(f"Generating with {style} notes...")
        generate_presentation(
            title=f"Artificial Intelligence Overview",
            topic=f"Understanding AI in Modern Technology",
            sections=test_sections,
            theme_name="Business Black and Yellow",
            notes_style=style,
            filename=f"test_AI_{style.replace(' ', '_')}.pptx"
        )
    
    print("\n✅ Test presentations generated with alternating layouts and LLM-style notes!")
