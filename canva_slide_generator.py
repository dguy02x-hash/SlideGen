"""
Slide Generator with Canva Themes and Image Placeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from canva_themes import get_canva_theme

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def create_slide_with_canva_theme(prs, theme_name, title, content, add_image_placeholder=True):
    """
    Create a slide with Canva theme applied and optional image placeholder
    
    Args:
        prs: Presentation object
        theme_name: Name of the Canva theme
        title: Slide title
        content: List of bullet points
        add_image_placeholder: Whether to add a blank image placeholder
    """
    theme = get_canva_theme(theme_name)
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    bg_rgb = hex_to_rgb(theme['background'])
    fill.fore_color.rgb = RGBColor(*bg_rgb)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(5.5) if add_image_placeholder else Inches(9)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    
    # Style title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = theme['title_font']
    title_para.font.size = Pt(theme['title_size'])
    title_para.font.bold = True
    title_rgb = hex_to_rgb(theme['title_color'])
    title_para.font.color.rgb = RGBColor(*title_rgb)
    
    # Add content
    content_left = Inches(0.5)
    content_top = Inches(1.8)
    content_width = Inches(5.5) if add_image_placeholder else Inches(9)
    content_height = Inches(4.5)
    
    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    # Add bullet points
    for i, bullet in enumerate(content):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = bullet
        p.level = 0
        p.font.name = theme['body_font']
        p.font.size = Pt(theme['body_size'])
        text_rgb = hex_to_rgb(theme['text_color'])
        p.font.color.rgb = RGBColor(*text_rgb)
        p.space_after = Pt(12)
    
    # Add image placeholder if requested
    if add_image_placeholder and theme['has_image_placeholder']:
        img_pos = theme['image_position']
        img_left = Inches(img_pos['left'])
        img_top = Inches(img_pos['top'])
        img_width = Inches(img_pos['width'])
        img_height = Inches(img_pos['height'])
        
        # Create a rectangle as image placeholder
        img_placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            img_left,
            img_top,
            img_width,
            img_height
        )
        
        # Style the placeholder
        img_placeholder.fill.solid()
        # Lighter version of background for placeholder
        accent_rgb = hex_to_rgb(theme['accent_color'])
        img_placeholder.fill.fore_color.rgb = RGBColor(*accent_rgb)
        img_placeholder.line.color.rgb = RGBColor(*text_rgb)
        img_placeholder.line.width = Pt(2)
        
        # Add "Image" text in center
        text_frame = img_placeholder.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "IMAGE"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(*text_rgb)
        text_frame.vertical_anchor = 1  # Middle
    
    return slide

def generate_presentation_with_canva_themes(title, topic, sections, theme_name='canva_navy_blue', add_images=True, filename='output.pptx'):
    """
    Generate a complete presentation with Canva theme
    
    Args:
        title: Presentation title
        topic: Presentation topic
        sections: List of sections with content
        theme_name: Which Canva theme to use
        add_images: Whether to add image placeholders
        filename: Output filename
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    theme = get_canva_theme(theme_name)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[6]  # Blank
    title_slide = prs.slides.add_slide(title_slide_layout)
    
    # Background
    background = title_slide.background
    fill = background.fill
    fill.solid()
    bg_rgb = hex_to_rgb(theme['background'])
    fill.fore_color.rgb = RGBColor(*bg_rgb)
    
    # Title
    title_box = title_slide.shapes.add_textbox(
        Inches(1),
        Inches(2.5),
        Inches(8),
        Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.name = theme['title_font']
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_rgb = hex_to_rgb(theme['title_color'])
    title_para.font.color.rgb = RGBColor(*title_rgb)
    
    # Content slides
    for section in sections:
        slide_title = section.get('title', 'Untitled')
        content = section.get('facts', [])
        
        create_slide_with_canva_theme(
            prs,
            theme_name,
            slide_title,
            content,
            add_image_placeholder=add_images
        )
    
    # Save
    prs.save(filename)
    return filename


# Example usage
if __name__ == '__main__':
    sections = [
        {
            'title': 'Introduction',
            'facts': [
                'Cloud computing has revolutionized IT infrastructure',
                'Organizations report 32% cost savings on average',
                'Enables global deployment across 25+ regions'
            ]
        },
        {
            'title': 'Key Benefits',
            'facts': [
                '99.99% uptime vs 95% traditional infrastructure',
                'Elastic scaling based on demand',
                'Reduced time-to-market by 60%'
            ]
        }
    ]
    
    generate_presentation_with_canva_themes(
        title='Cloud Computing Overview',
        topic='Technology',
        sections=sections,
        theme_name='canva_navy_blue',
        add_images=True,
        filename='test_canva_presentation.pptx'
    )
    print("Presentation generated!")
